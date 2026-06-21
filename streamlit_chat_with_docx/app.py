import streamlit as st
import streamlit.components.v1 as components
import seaborn as sns
import os
from openai import OpenAI
from dotenv import load_dotenv
from datetime import datetime
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import io

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

load_dotenv()


# ---------------- PAGE CONFIG ----------------

st.set_page_config(
    page_title="AI Document & Data Assistant 🤖📊",
    page_icon="📚",
    layout="wide"
)


# ---------------- STYLE ----------------

st.markdown("""
<style>
.main { background-color: #f8f9fb; }
.stChatMessage { border-radius: 12px; padding: 12px; }
.upload-box { background-color: #eef6ff; padding: 20px; border-radius: 10px; }
.title { font-size: 38px; font-weight: bold; text-align:center; }
</style>
""", unsafe_allow_html=True)


# ---------------- SESSION STATE ----------------

if "doc" not in st.session_state:
    st.session_state.doc = None

if "name" not in st.session_state:
    st.session_state.name = None

if "messages" not in st.session_state:
    st.session_state.messages = []

if "df" not in st.session_state:
    st.session_state.df = None


# ---------------- DOCUMENT READER ----------------

def extract_text(file):

    name = file.name.lower()

    if name.endswith(".txt"):
        return file.read().decode("utf-8")

    elif name.endswith(".pdf"):
        reader = PdfReader(file)
        return "\n".join([p.extract_text() or "" for p in reader.pages])

    elif name.endswith(".docx"):
        doc = Document(file)
        return "\n".join([p.text for p in doc.paragraphs])

    elif name.endswith(".pptx"):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif name.endswith(".xlsx"):
        df = pd.read_excel(file)
        return df.to_string()

    elif name.endswith(".csv"):
        df = pd.read_csv(file)
        return df.to_string()

    else:
        return "Unsupported file type"


# ---------------- AI FOR DOCUMENT CHAT ----------------

def ask_ai(messages, document):

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    system = f"""
You are a smart assistant 🤖.

Answer ONLY from document.

Document:
{document}

Rules:
- Be short and clear
- Use bullets if needed
- If not found say: "Not found in document"
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": system}, *messages],
        temperature=0.3
    )

    return response.choices[0].message.content


# ---------------- AI FOR DATA ANALYSIS ----------------

def generate_code(question, columns):

    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

    prompt = f"""
You are a data analyst 📊.

DataFrame name = df

Columns:
{columns}

User request:
{question}

Rules:
- use pandas, matplotlib, seaborn
- DO NOT use plt.show()
- save chart in variable fig OR directly use plt
- keep code simple and correct
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2
    )

    return response.choices[0].message.content


# ---------------- SIDEBAR ----------------

with st.sidebar:

    st.title("🤖 AI Assistant")

    api = st.text_input("🔑 OpenAI API Key", type="password")

    if api:
        os.environ["OPENAI_API_KEY"] = api

    st.divider()

    if st.button("🗑️ Clear Chat"):
        st.session_state.messages = []
        st.rerun()

    st.info("📌 Upload documents or CSV to start")


# ---------------- TITLE ----------------

st.markdown('<div class="title">📚 AI Document + Data Analyst 📊</div>', unsafe_allow_html=True)


# ---------------- FILE UPLOAD ----------------

st.markdown("## 📄 Upload Document")

uploaded = st.file_uploader(
    "Upload PDF, DOCX, TXT, PPTX, Excel, CSV",
    type=["pdf", "docx", "txt", "pptx", "xlsx", "csv"]
)


# ---------------- PROCESS FILE ----------------

if uploaded:

    if uploaded.name.endswith(".csv"):

        st.session_state.df = pd.read_csv(uploaded)

        st.success("CSV Loaded 📊")

        st.dataframe(st.session_state.df.head())

    else:

        st.session_state.doc = extract_text(uploaded)
        st.session_state.name = uploaded.name

        st.success(f"{uploaded.name} loaded 📄")


# ---------------- CSV ANALYSIS MODE ----------------

if st.session_state.df is not None:

    st.subheader("📊 Ask Questions on CSV Data")

    q = st.text_input("Ask about dataset (e.g., plot sales trend)")

    if q:

        with st.spinner("Generating Python code 🤖"):

            code = generate_code(q, str(st.session_state.df.columns.tolist()))

        st.code(code, language="python")

        local = {
            "df": st.session_state.df,
            "pd": pd,
            "plt": plt,
            "sns": sns
        }

        try:
            exec(code, local)
            st.pyplot(plt)

        except Exception as e:
            st.error(f"Error: {e}")


# ---------------- DOCUMENT CHAT MODE ----------------

elif st.session_state.doc is not None:

    st.subheader(f"📄 Chat with {st.session_state.name}")

    for m in st.session_state.messages:
        with st.chat_message(m["role"]):
            st.write(m["content"])

    user = st.chat_input("Ask anything about document...")

    if user:

        st.session_state.messages.append({"role": "user", "content": user})

        with st.chat_message("user"):
            st.write(user)

        with st.chat_message("assistant"):

            with st.spinner("Thinking 🤖"):

                answer = ask_ai(
                    st.session_state.messages,
                    st.session_state.doc
                )

                st.write(answer)

        st.session_state.messages.append(
            {"role": "assistant", "content": answer}
        )

else:

    st.info("⬆️ Upload a document or CSV to start")
